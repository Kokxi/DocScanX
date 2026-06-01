"""文件解析器 — 按格式分派解析器，提取纯文本。"""
import os
import json
import xml.etree.ElementTree as ET
from dataclasses import dataclass
from typing import Optional

import chardet


@dataclass
class ParseResult:
    text: str
    encoding: str = "utf-8"
    page_count: int = 1
    error: Optional[str] = None
    metadata: Optional[dict] = None


def _read_text_file(file_path: str) -> ParseResult:
    """读取纯文本文件，自动检测编码。"""
    with open(file_path, "rb") as f:
        raw = f.read()
    if not raw:
        return ParseResult(text="", encoding="utf-8")

    detected = chardet.detect(raw)
    encoding = detected.get("encoding", "utf-8") or "utf-8"

    try:
        text = raw.decode(encoding)
    except (UnicodeDecodeError, LookupError):
        text = raw.decode("latin-1")
        encoding = "latin-1"

    return ParseResult(text=text, encoding=encoding)


def _parse_docx(file_path: str) -> ParseResult:
    try:
        from docx import Document
        doc = Document(file_path)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells if cell.text)
                if row_text.strip():
                    parts.append(row_text)
        return ParseResult(text="\n".join(parts))
    except Exception as e:
        return ParseResult(text="", error=f"docx 解析失败: {e}")


def _parse_xlsx(file_path: str) -> ParseResult:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"[Sheet: {sheet_name}]")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    parts.append(row_text)
        wb.close()
        return ParseResult(text="\n".join(parts))
    except Exception as e:
        return ParseResult(text="", error=f"xlsx 解析失败: {e}")


def _parse_pptx(file_path: str) -> ParseResult:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_texts.append(para.text)
            if slide_texts:
                parts.append(f"[Slide {i}]")
                parts.extend(slide_texts)
        return ParseResult(text="\n".join(parts), page_count=len(prs.slides))
    except Exception as e:
        return ParseResult(text="", error=f"pptx 解析失败: {e}")


def _parse_pdf(file_path: str) -> ParseResult:
    try:
        import pdfplumber
        parts = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    parts.append(text)
            page_count = len(pdf.pages)
        return ParseResult(text="\n".join(parts), page_count=page_count)
    except Exception as e:
        return ParseResult(text="", error=f"pdf 解析失败: {e}")


def _parse_json_file(file_path: str) -> ParseResult:
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        text = json.dumps(data, ensure_ascii=False, indent=2)
        return ParseResult(text=text)
    except Exception as e:
        return ParseResult(text="", error=f"json 解析失败: {e}")


def _parse_xml_file(file_path: str) -> ParseResult:
    try:
        tree = ET.parse(file_path)
        parts = []

        def walk(elem, depth=0):
            if elem.text and elem.text.strip():
                parts.append(f"{'  ' * depth}{elem.tag}: {elem.text.strip()}")
            for child in elem:
                walk(child, depth + 1)

        walk(tree.getroot())
        return ParseResult(text="\n".join(parts))
    except Exception as e:
        return ParseResult(text="", error=f"xml 解析失败: {e}")


PARSERS = {
    ".txt": _read_text_file,
    ".csv": _read_text_file,
    ".md": _read_text_file,
    ".log": _read_text_file,
    ".docx": _parse_docx,
    ".xlsx": _parse_xlsx,
    ".pptx": _parse_pptx,
    ".pdf": _parse_pdf,
    ".json": _parse_json_file,
    ".xml": _parse_xml_file,
    ".py": _read_text_file,
    ".java": _read_text_file,
    ".js": _read_text_file,
    ".html": _read_text_file,
    ".sql": _read_text_file,
    ".css": _read_text_file,
    ".cpp": _read_text_file,
}

UNSUPPORTED_EXTS = {".doc", ".xls", ".ppt"}


def parse_file(file_path: str) -> ParseResult:
    """解析单个文件，返回纯文本内容。"""
    ext = os.path.splitext(file_path)[1].lower()

    if ext in UNSUPPORTED_EXTS:
        return ParseResult(text="", error=f"老旧格式暂不支持: {ext}")

    parser = PARSERS.get(ext)
    if parser is None:
        return ParseResult(text="", error=f"不支持的文件格式: {ext}")

    try:
        return parser(file_path)
    except Exception as e:
        return ParseResult(text="", error=str(e))
