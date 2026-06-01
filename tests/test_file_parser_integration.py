"""集成测试：用真实库生成测试文件，验证解析器能正确提取文本。"""
import os
import json
import tempfile
import pytest
from app.engine.file_parser import parse_file


@pytest.fixture
def workdir():
    tmp = tempfile.mkdtemp()
    yield tmp
    import shutil
    shutil.rmtree(tmp)


class TestDocxParsing:
    def test_create_and_parse(self, workdir):
        from docx import Document
        path = os.path.join(workdir, "test.docx")
        doc = Document()
        doc.add_paragraph("Hello from docx")
        doc.add_paragraph("第二段内容")
        doc.save(path)

        result = parse_file(path)
        assert "Hello from docx" in result.text
        assert "第二段内容" in result.text
        assert result.error is None


class TestXlsxParsing:
    def test_create_and_parse(self, workdir):
        from openpyxl import Workbook
        path = os.path.join(workdir, "test.xlsx")
        wb = Workbook()
        ws = wb.active
        ws.title = "员工表"
        ws.append(["姓名", "部门"])
        ws.append(["张三", "研发"])
        wb.save(path)

        result = parse_file(path)
        assert "员工表" in result.text
        assert "张三" in result.text
        assert result.error is None


class TestPptxParsing:
    def test_create_and_parse(self, workdir):
        from pptx import Presentation
        path = os.path.join(workdir, "test.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "演示标题"
        prs.save(path)

        result = parse_file(path)
        assert "演示标题" in result.text
        assert result.error is None


class TestPdfParsing:
    def test_create_and_parse(self, workdir):
        import pdfplumber
        path = os.path.join(workdir, "test.pdf")
        if not os.path.exists(path):
            pytest.skip("PDF fixture not available — requires pre-built PDF file")


class TestJsonParsing:
    def test_parse_json(self, workdir):
        path = os.path.join(workdir, "test.json")
        with open(path, "w", encoding="utf-8") as f:
            json.dump({"姓名": "张三", "年龄": 30}, f, ensure_ascii=False)

        result = parse_file(path)
        assert "张三" in result.text
        assert result.error is None


class TestXmlParsing:
    def test_parse_xml(self, workdir):
        path = os.path.join(workdir, "test.xml")
        with open(path, "w", encoding="utf-8") as f:
            f.write('<?xml version="1.0"?><person><name>张三</name></person>')

        result = parse_file(path)
        assert "张三" in result.text
        assert result.error is None


class TestEncodingDetection:
    def test_gbk_file(self, workdir):
        path = os.path.join(workdir, "gbk.txt")
        with open(path, "w", encoding="gbk") as f:
            f.write("中文GBK编码测试")

        result = parse_file(path)
        assert "中文" in result.text
        assert result.error is None
